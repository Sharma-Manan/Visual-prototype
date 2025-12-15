from pptx import Presentation
from pptx.util import Inches

def create_pptx(slide_data, out_path, placeholder_path="assets/icon_default.png"):
    """
    Builds a PowerPoint presentation from structured slide data.
    
    Args:
        slide_data (list): List of dictionaries containing slide content.
        out_path (str): Path to save the output PPTX file.
        placeholder_path (str): Path to the default image/icon.
        
    Returns:
        str: Path to the saved PPTX file.
    """
    prs = Presentation()
    
    # Use layout index 1 (Title and Content)
    slide_layout = prs.slide_layouts[1]
    
    for data in slide_data:
        slide = prs.slides.add_slide(slide_layout)
        
        # 1. Title
        title = slide.shapes.title
        title.text = data['title']
        
        # 2. Image Placeholder (top-right area)
        try:
            left = Inches(5.5)
            top = Inches(1.5)
            width = height = Inches(3.0)
            slide.shapes.add_picture(placeholder_path, left, top, width, height)
        except FileNotFoundError:
            print("    - WARNING: Placeholder image not found. Skipping image.")

        # 3. Text/Bullet Placeholder (left area)
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear()
        
        for bullet in data['bullets']:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            
    prs.save(out_path)
    return out_path